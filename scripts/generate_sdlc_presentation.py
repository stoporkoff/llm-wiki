from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "docs" / "assets"
OUTPUT = ROOT / "docs" / "SDLC-brainstorm.pptx"

W = Inches(13.333)
H = Inches(7.5)
INK = RGBColor(20, 27, 33)
MUTED = RGBColor(83, 96, 107)
PAPER = RGBColor(250, 249, 246)
WHITE = RGBColor(255, 255, 255)
TEAL = RGBColor(56, 182, 184)
CYAN = RGBColor(91, 194, 225)
ORANGE = RGBColor(255, 166, 67)
CORAL = RGBColor(244, 107, 77)
PURPLE = RGBColor(166, 120, 193)
BLUE = RGBColor(109, 147, 220)
GREEN = RGBColor(134, 190, 122)
PALE_BLUE = RGBColor(229, 240, 246)
PALE_TEAL = RGBColor(225, 245, 242)
PALE_ORANGE = RGBColor(255, 239, 214)
PALE_PURPLE = RGBColor(241, 232, 246)
PALE_GREEN = RGBColor(232, 245, 226)


def set_background(slide, color: RGBColor = PAPER) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: int = 18,
    bold: bool = False,
    color: RGBColor = INK,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    font: str = "Aptos",
    margin: float = 0.04,
) -> object:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    run = paragraph.runs[0]
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_title(slide, title: str, eyebrow: str | None = None) -> None:
    if eyebrow:
        add_text(slide, eyebrow.upper(), 0.7, 0.32, 5.8, 0.28, size=10, bold=True, color=TEAL)
    add_text(slide, title, 0.67, 0.62, 12.0, 0.72, size=27, bold=True)
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.34), Inches(11.9), Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(205, 213, 219)
    line.line.fill.background()


def add_footer(slide, number: int, section: str = "SDLC / Brainstorm") -> None:
    add_text(slide, section, 0.7, 7.08, 4.0, 0.18, size=8, color=MUTED)
    add_text(
        slide,
        f"{number:02d}",
        12.05,
        7.04,
        0.55,
        0.22,
        size=9,
        bold=True,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )


def add_card(
    slide,
    title: str,
    body: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    accent: RGBColor = TEAL,
    fill: RGBColor = WHITE,
    title_size: int = 16,
    body_size: int = 12,
) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = RGBColor(205, 213, 219)
    shape.line.width = Pt(1.1)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.09), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    add_text(slide, title, x + 0.25, y + 0.2, w - 0.45, 0.4, size=title_size, bold=True)
    add_text(slide, body, x + 0.25, y + 0.72, w - 0.45, h - 0.86, size=body_size, color=MUTED)


def add_pill(slide, text: str, x: float, y: float, w: float, color: RGBColor) -> None:
    pill = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.42)
    )
    pill.fill.solid()
    pill.fill.fore_color.rgb = color
    pill.line.color.rgb = INK
    pill.line.width = Pt(0.8)
    frame = pill.text_frame
    frame.clear()
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = frame.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    r = p.runs[0]
    r.font.name = "Aptos"
    r.font.size = Pt(10)
    r.font.bold = True
    r.font.color.rgb = INK


def add_chevron(slide, x: float, y: float, color: RGBColor = TEAL) -> None:
    c = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Inches(x), Inches(y), Inches(0.34), Inches(0.5))
    c.fill.solid()
    c.fill.fore_color.rgb = color
    c.line.fill.background()


def add_bullets(
    slide, items: list[str], x: float, y: float, w: float, h: float, size: int = 15
) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for index, item in enumerate(items):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(11)
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = INK
        p.text = f"•  {item}"


def new_slide(prs: Presentation, number: int, title: str, eyebrow: str) -> object:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    add_title(slide, title, eyebrow)
    add_footer(slide, number)
    return slide


def build() -> None:
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_background(slide)
    slide.shapes.add_picture(str(ASSETS / "sdlc-factory-hero.png"), 0, Inches(1.18), width=W)
    veil = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(1.75))
    veil.fill.solid()
    veil.fill.fore_color.rgb = PAPER
    veil.line.fill.background()
    add_text(slide, "SDLC", 0.7, 0.34, 2.2, 0.7, size=34, bold=True)
    add_text(
        slide,
        "AI-assisted software delivery control plane",
        2.5,
        0.48,
        8.8,
        0.5,
        size=22,
        bold=True,
    )
    add_text(
        slide,
        "Brainstorm deck · operating model, evidence, gates, and next bets",
        2.52,
        1.03,
        8.4,
        0.32,
        size=12,
        color=MUTED,
    )
    add_pill(slide, "DISCOVER", 0.72, 6.72, 1.25, PALE_TEAL)
    add_pill(slide, "BUILD", 2.08, 6.72, 1.05, PALE_ORANGE)
    add_pill(slide, "VERIFY", 3.24, 6.72, 1.15, PALE_PURPLE)
    add_pill(slide, "DELIVER", 4.5, 6.72, 1.2, PALE_GREEN)
    add_footer(slide, 1)

    slide = new_slide(prs, 2, "Why SDLC exists", "Problem framing")
    add_card(
        slide,
        "Speed without evidence",
        "Agents can generate quickly, but a fast answer is not automatically a usable delivery.",
        0.75,
        1.72,
        3.75,
        1.45,
        accent=ORANGE,
        fill=PALE_ORANGE,
    )
    add_card(
        slide,
        "Specialists without control",
        "Parallel roles increase throughput only when ownership, tools, and handoffs are explicit.",
        4.8,
        1.72,
        3.75,
        1.45,
        accent=PURPLE,
        fill=PALE_PURPLE,
    )
    add_card(
        slide,
        "Automation without trust",
        "A workflow needs deterministic gates, runtime proof, provenance, "
        "and a fail-closed release decision.",
        8.85,
        1.72,
        3.75,
        1.45,
        accent=TEAL,
        fill=PALE_TEAL,
    )
    add_text(slide, "The product thesis", 0.8, 3.72, 3.0, 0.4, size=20, bold=True)
    add_text(
        slide,
        "Turn an ambiguous goal into an observable, reviewable, locally runnable "
        "artifact — with every decision and gate recorded.",
        0.8,
        4.28,
        11.4,
        1.0,
        size=24,
        bold=True,
    )
    add_pill(slide, "LOCAL-FIRST", 0.82, 5.75, 1.55, PALE_BLUE)
    add_pill(slide, "FAIL-CLOSED", 2.55, 5.75, 1.62, PALE_ORANGE)
    add_pill(slide, "EVIDENCE-LINKED", 4.35, 5.75, 1.9, PALE_TEAL)
    add_pill(slide, "REUSABLE", 6.43, 5.75, 1.35, PALE_PURPLE)
    add_pill(slide, "OBSERVABLE", 7.96, 5.75, 1.55, PALE_GREEN)

    slide = new_slide(prs, 3, "One goal, a governed delivery system", "Operating model")
    stages = [
        ("Discovery", "Bound the goal\nand acceptance", TEAL, PALE_TEAL),
        ("Planning", "Architecture +\nownership plan", BLUE, PALE_BLUE),
        ("Build", "Specialists write\nisolated areas", ORANGE, PALE_ORANGE),
        ("Verify", "Runtime tests +\nsecurity review", PURPLE, PALE_PURPLE),
        ("Review", "Cross-role trust\nverdict", CORAL, RGBColor(252, 228, 221)),
        ("Package", "Score + reusable\ntool contract", GREEN, PALE_GREEN),
    ]
    x = 0.63
    for index, (name, body, accent, fill) in enumerate(stages):
        add_card(
            slide,
            name,
            body,
            x,
            2.0,
            1.75,
            2.0,
            accent=accent,
            fill=fill,
            title_size=14,
            body_size=11,
        )
        if index < len(stages) - 1:
            add_chevron(slide, x + 1.83, 2.74, accent)
        x += 2.08
    add_text(
        slide,
        "Every stage emits events, agent outputs, token/tool telemetry, "
        "and a durable session artifact.",
        1.15,
        4.65,
        11.0,
        0.62,
        size=18,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_pill(slide, "SQLite session state", 2.15, 5.55, 1.8, PALE_BLUE)
    add_pill(slide, "session.yaml", 4.15, 5.55, 1.45, PALE_ORANGE)
    add_pill(slide, "OpenTelemetry", 5.8, 5.55, 1.65, PALE_PURPLE)
    add_pill(slide, "local preview", 7.65, 5.55, 1.55, PALE_GREEN)

    slide = new_slide(prs, 4, "Practices encoded as product behavior", "Engineering system")
    practices = [
        (
            "Explicit FSM",
            "No hidden phase changes; terminal states are durable and inspectable.",
            TEAL,
        ),
        (
            "Bounded agency",
            "Each role has an allowlisted toolset and workspace write boundary.",
            BLUE,
        ),
        (
            "Parallel where safe",
            "Build and verification concurrency without overlapping ownership.",
            PURPLE,
        ),
        (
            "Runtime evidence",
            "Test commands and outputs are first-class gate inputs, not prose decoration.",
            ORANGE,
        ),
        (
            "Defensive review",
            "Security and final review are read-only and cannot patch around findings.",
            CORAL,
        ),
        (
            "Reusable delivery",
            "Approved work is scored, versioned, and exposed through a strict contract.",
            GREEN,
        ),
    ]
    for i, (name, body, accent) in enumerate(practices):
        col, row = i % 3, i // 3
        add_card(slide, name, body, 0.72 + col * 4.15, 1.67 + row * 2.35, 3.82, 1.9, accent=accent)

    slide = new_slide(prs, 5, "Quality gates must fail closed", "Trust model")
    slide.shapes.add_picture(
        str(ASSETS / "sdlc-quality-gates.png"), Inches(0.35), Inches(1.32), width=Inches(12.65)
    )
    label = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.1), Inches(5.97), Inches(4.45), Inches(0.73)
    )
    label.fill.solid()
    label.fill.fore_color.rgb = WHITE
    label.line.color.rgb = CORAL
    label.line.width = Pt(1.4)
    add_text(
        slide,
        "BLOCKED means stop here — never “passed with caveats.”",
        8.3,
        6.13,
        4.0,
        0.35,
        size=13,
        bold=True,
        color=CORAL,
        align=PP_ALIGN.CENTER,
    )

    slide = new_slide(prs, 6, "What happened in the latest run", "Incident review · 19 Aug 2026")
    events = [
        ("23:08", "Session started", "React login UI requested", TEAL),
        ("23:10", "Implementation", "Static requirements satisfied", BLUE),
        ("23:11", "QA blocked", "vitest: not found · exit 127", ORANGE),
        ("23:11", "False pass", "Verification marked passed", CORAL),
        ("23:12", "Final rejection", "Runtime evidence missing", PURPLE),
    ]
    x = 0.78
    for i, (time, name, body, accent) in enumerate(events):
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x), Inches(2.05), Inches(0.55), Inches(0.55)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = accent
        circle.line.color.rgb = INK
        if i < len(events) - 1:
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(x + 0.55), Inches(2.29), Inches(1.65), Inches(0.08)
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = RGBColor(192, 203, 210)
            bar.line.fill.background()
        add_text(
            slide,
            time,
            x - 0.05,
            1.65,
            0.7,
            0.25,
            size=10,
            bold=True,
            color=accent,
            align=PP_ALIGN.CENTER,
        )
        add_text(slide, name, x - 0.3, 2.82, 1.25, 0.42, size=13, bold=True, align=PP_ALIGN.CENTER)
        add_text(
            slide, body, x - 0.52, 3.35, 1.7, 0.82, size=10, color=MUTED, align=PP_ALIGN.CENTER
        )
        x += 2.45
    add_card(
        slide,
        "Primary cause",
        "Dependency setup happened after QA, during Preview. QA therefore could not "
        "run the generated frontend suite.",
        0.82,
        4.75,
        3.75,
        1.25,
        accent=ORANGE,
        fill=PALE_ORANGE,
        body_size=11,
    )
    add_card(
        slide,
        "Control failure",
        "The Verification gate accepted narrative output and did not interpret "
        "QA BLOCKED as terminal.",
        4.8,
        4.75,
        3.75,
        1.25,
        accent=CORAL,
        fill=RGBColor(252, 228, 221),
        body_size=11,
    )
    add_card(
        slide,
        "Now fixed",
        "QA prepares a lockfile-based environment; machine-readable verdicts stop "
        "the workflow before review.",
        8.78,
        4.75,
        3.75,
        1.25,
        accent=GREEN,
        fill=PALE_GREEN,
        body_size=11,
    )

    slide = new_slide(prs, 7, "The fix: deterministic tests + explicit verdicts", "Control design")
    add_card(
        slide,
        "1 · Dependency policy",
        "Reject `latest`. Generate and preserve package-lock.json when absent.",
        0.75,
        1.72,
        3.7,
        1.4,
        accent=BLUE,
        fill=PALE_BLUE,
    )
    add_card(
        slide,
        "2 · Reproducible execution",
        "Install with npm ci --ignore-scripts, run the suite, then remove node_modules.",
        4.82,
        1.72,
        3.7,
        1.4,
        accent=TEAL,
        fill=PALE_TEAL,
    )
    add_card(
        slide,
        "3 · Gate contract",
        "Require QA PASSED and SECURITY PASSED as the first non-empty line.",
        8.88,
        1.72,
        3.7,
        1.4,
        accent=PURPLE,
        fill=PALE_PURPLE,
    )
    add_text(
        slide,
        "QA BLOCKED",
        1.08,
        4.15,
        2.2,
        0.52,
        size=22,
        bold=True,
        color=CORAL,
        align=PP_ALIGN.CENTER,
    )
    add_chevron(slide, 3.35, 4.16, CORAL)
    add_text(
        slide,
        "Verification raises",
        3.9,
        4.15,
        2.3,
        0.52,
        size=18,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_chevron(slide, 6.35, 4.16, CORAL)
    add_text(
        slide,
        "Session → FAILED",
        6.9,
        4.15,
        2.4,
        0.52,
        size=22,
        bold=True,
        color=CORAL,
        align=PP_ALIGN.CENTER,
    )
    add_chevron(slide, 9.42, 4.16, CORAL)
    add_text(
        slide,
        "Reviewer not invoked",
        9.95,
        4.15,
        2.2,
        0.52,
        size=17,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "Regression test: blocked QA cannot emit stage-completed and cannot start final review.",
        1.2,
        5.55,
        10.9,
        0.55,
        size=18,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    slide = new_slide(prs, 8, "System architecture", "Current platform")
    add_card(
        slide,
        "Factory UI + API",
        "FastAPI control plane\nSession and event APIs\nLocal preview routing",
        0.75,
        1.68,
        2.65,
        2.0,
        accent=TEAL,
        fill=PALE_TEAL,
    )
    add_chevron(slide, 3.52, 2.45, TEAL)
    add_card(
        slide,
        "Orchestrator",
        "FSM + stage pipeline\nAgent gateway\nDeterministic gates",
        3.98,
        1.68,
        2.65,
        2.0,
        accent=BLUE,
        fill=PALE_BLUE,
    )
    add_chevron(slide, 6.75, 2.45, BLUE)
    add_card(
        slide,
        "Session workspace",
        "Role-owned files\nTests + deployment\nSession artifact",
        7.2,
        1.68,
        2.65,
        2.0,
        accent=ORANGE,
        fill=PALE_ORANGE,
    )
    add_chevron(slide, 9.97, 2.45, ORANGE)
    add_card(
        slide,
        "Delivery",
        "Exportable project\nReusable tool\nVerified preview",
        10.42,
        1.68,
        2.15,
        2.0,
        accent=GREEN,
        fill=PALE_GREEN,
    )
    add_pill(slide, "OpenAI Responses", 1.0, 4.52, 1.9, PALE_PURPLE)
    add_pill(slide, "SQLite", 3.2, 4.52, 1.2, PALE_BLUE)
    add_pill(slide, "File wiki", 4.7, 4.52, 1.25, PALE_TEAL)
    add_pill(slide, "OTel Collector", 6.25, 4.52, 1.65, PALE_ORANGE)
    add_pill(slide, "Jaeger", 8.2, 4.52, 1.1, PALE_PURPLE)
    add_pill(slide, "Prometheus", 9.6, 4.52, 1.45, PALE_GREEN)
    add_pill(slide, "Grafana", 11.35, 4.52, 1.05, PALE_BLUE)
    add_text(
        slide,
        "Local Docker Compose keeps the delivery loop inspectable: code, traces, "
        "metrics, events, and artifacts share one operational boundary.",
        1.05,
        5.55,
        11.2,
        0.72,
        size=17,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    slide = new_slide(prs, 9, "What we should measure", "Product + engineering metrics")
    metrics = [
        (
            "Flow",
            "Lead time per stage\nQueue and agent duration\nParallelism efficiency",
            TEAL,
            PALE_TEAL,
        ),
        (
            "Quality",
            "First-pass acceptance\nGate failure reasons\nEscaped defects",
            ORANGE,
            PALE_ORANGE,
        ),
        (
            "Trust",
            "Runtime evidence coverage\nCritical findings\nHuman override rate",
            PURPLE,
            PALE_PURPLE,
        ),
        (
            "Economics",
            "Tokens per accepted delivery\nTool reuse savings\nCost by workflow",
            GREEN,
            PALE_GREEN,
        ),
    ]
    for i, (name, body, accent, fill) in enumerate(metrics):
        add_card(
            slide,
            name,
            body,
            0.75 + i * 3.1,
            1.75,
            2.75,
            2.15,
            accent=accent,
            fill=fill,
            title_size=18,
            body_size=12,
        )
    add_text(slide, "North-star candidate", 0.85, 4.65, 2.3, 0.36, size=13, bold=True, color=TEAL)
    add_text(
        slide,
        "Verified, locally runnable deliveries per unit of time and cost",
        0.85,
        5.12,
        11.5,
        0.75,
        size=25,
        bold=True,
    )
    add_text(
        slide,
        "Guardrail: the metric counts only deliveries that passed runtime QA, "
        "security, final review, packaging, and preview.",
        0.87,
        6.05,
        10.8,
        0.42,
        size=13,
        color=MUTED,
    )

    slide = new_slide(prs, 10, "Brainstorm: decisions worth making tomorrow", "Workshop prompts")
    questions = [
        "Which workflows deserve hard deterministic gates versus human review?",
        "What is the minimum runtime evidence for frontend, API, data, and infrastructure work?",
        "Should agents repair a blocked delivery automatically, or open a bounded "
        "remediation loop?",
        "How do trusted reusable tools earn, retain, and lose trust over time?",
        "Where should policy live: prompts, code, signed manifests, or all three?",
        "Which telemetry is useful without turning agent work into vanity dashboards?",
    ]
    for i, q in enumerate(questions):
        col, row = i % 2, i // 2
        accent = [TEAL, ORANGE, PURPLE, BLUE, CORAL, GREEN][i]
        add_card(
            slide,
            f"0{i + 1}",
            q,
            0.78 + col * 6.15,
            1.62 + row * 1.65,
            5.72,
            1.3,
            accent=accent,
            title_size=13,
            body_size=12,
        )

    slide = new_slide(prs, 11, "A pragmatic 30 / 60 / 90-day path", "Roadmap hypothesis")
    add_card(
        slide,
        "30 days · Make it reliable",
        "Structured verdict schema\nDependency + browser test harness\n"
        "Failure taxonomy\nGolden workflow fixtures",
        0.78,
        1.75,
        3.7,
        3.65,
        accent=TEAL,
        fill=PALE_TEAL,
        title_size=18,
        body_size=14,
    )
    add_card(
        slide,
        "60 days · Make it adaptive",
        "Bounded remediation loops\nPolicy profiles by project type\n"
        "Artifact provenance + SBOM\nReusable tool lifecycle",
        4.82,
        1.75,
        3.7,
        3.65,
        accent=PURPLE,
        fill=PALE_PURPLE,
        title_size=18,
        body_size=14,
    )
    add_card(
        slide,
        "90 days · Make it learn",
        "Outcome-linked scoring\nCross-run pattern analysis\n"
        "Human feedback calibration\nPortfolio-level controls",
        8.86,
        1.75,
        3.7,
        3.65,
        accent=ORANGE,
        fill=PALE_ORANGE,
        title_size=18,
        body_size=14,
    )
    add_text(
        slide,
        "Constraint throughout: local-first, auditable, reversible, and fail-closed.",
        1.2,
        6.02,
        10.9,
        0.45,
        size=18,
        bold=True,
        align=PP_ALIGN.CENTER,
    )

    slide = new_slide(prs, 12, "Tomorrow’s desired outcome", "Close")
    add_text(
        slide,
        "Agree on the control plane we want — not just the agents we can run.",
        0.85,
        1.75,
        11.7,
        0.78,
        size=30,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_card(
        slide,
        "1 · Scope",
        "Pick one or two delivery archetypes for the next iteration.",
        1.05,
        3.05,
        3.35,
        1.45,
        accent=TEAL,
        fill=PALE_TEAL,
    )
    add_card(
        slide,
        "2 · Evidence",
        "Define the runtime proof that makes each archetype trustworthy.",
        4.98,
        3.05,
        3.35,
        1.45,
        accent=PURPLE,
        fill=PALE_PURPLE,
    )
    add_card(
        slide,
        "3 · Learning loop",
        "Decide how failures become policy, tests, and reusable assets.",
        8.9,
        3.05,
        3.35,
        1.45,
        accent=ORANGE,
        fill=PALE_ORANGE,
    )
    add_text(slide, "SDLC", 5.2, 5.35, 2.9, 0.7, size=34, bold=True, align=PP_ALIGN.CENTER)
    add_text(
        slide,
        "Software delivery with evidence at every handoff",
        3.25,
        6.08,
        6.85,
        0.42,
        size=16,
        bold=True,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    reopened = Presentation(OUTPUT)
    if len(reopened.slides) != 12 or OUTPUT.stat().st_size < 1_000_000:
        raise RuntimeError("Generated presentation failed structural validation")
    print(f"Wrote {OUTPUT} ({len(reopened.slides)} slides, {OUTPUT.stat().st_size} bytes)")


if __name__ == "__main__":
    build()
