import csv
import json
from io import BytesIO, StringIO
from pathlib import Path
from typing import Any


class UnsupportedFormatError(ValueError):
    pass


def extract_text(path: Path) -> str:
    if not path.is_file():
        raise ValueError(f"Source is not a file: {path}")
    suffix = path.suffix.casefold()
    if suffix in {".md", ".txt", ".rst", ".xml", ".yaml", ".yml"}:
        return path.read_text(encoding="utf-8-sig").strip()
    if suffix == ".json":
        value = json.loads(path.read_text(encoding="utf-8-sig"))
        return json.dumps(value, indent=2, ensure_ascii=False)
    if suffix in {".csv", ".tsv"}:
        return _extract_delimited(path, "\t" if suffix == ".tsv" else ",")
    if suffix == ".pdf":
        return _extract_pdf(path)
    if suffix == ".docx":
        return _extract_docx(path)
    if suffix == ".pptx":
        return _extract_pptx(path)
    if suffix == ".xlsx":
        return _extract_xlsx(path)
    raise UnsupportedFormatError(f"Unsupported source format: {suffix or '<none>'}")


def _extract_delimited(path: Path, delimiter: str) -> str:
    content = path.read_text(encoding="utf-8-sig")
    rows = list(csv.reader(StringIO(content), delimiter=delimiter))
    return "\n".join(" | ".join(cell.strip() for cell in row) for row in rows)


def _extract_pdf(path: Path) -> str:
    from pypdf import PdfReader

    reader = PdfReader(BytesIO(path.read_bytes()))
    pages = []
    for index, page in enumerate(reader.pages, start=1):
        pages.append(f"## Page {index}\n\n{(page.extract_text() or '').strip()}")
    return "\n\n".join(pages)


def _extract_docx(path: Path) -> str:
    from docx import Document

    document = Document(BytesIO(path.read_bytes()))
    paragraphs = (
        paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text
    )
    return "\n\n".join(paragraphs)


def _extract_pptx(path: Path) -> str:
    from pptx import Presentation

    presentation = Presentation(BytesIO(path.read_bytes()))
    slides = []
    for index, slide in enumerate(presentation.slides, start=1):
        text = [
            shape.text.strip()
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text
        ]
        slides.append(f"## Slide {index}\n\n" + "\n\n".join(text))
    return "\n\n".join(slides)


def _extract_xlsx(path: Path) -> str:
    from openpyxl import load_workbook

    workbook = load_workbook(path, read_only=True, data_only=True)
    sheets = []
    for worksheet in workbook.worksheets:
        rows = []
        for row in worksheet.iter_rows(values_only=True):
            values = [_stringify(value) for value in row]
            if any(values):
                rows.append(" | ".join(values))
        sheets.append(f"## Sheet: {worksheet.title}\n\n" + "\n".join(rows))
    workbook.close()
    return "\n\n".join(sheets)


def _stringify(value: Any) -> str:
    return "" if value is None else str(value).strip()
